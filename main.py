from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
from pptx import Presentation

app = FastAPI()

@app.post("/extract-text/")
async def extract_text(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        return JSONResponse(status_code=400, content={"error": "Only .pptx files are supported."})

    contents = await file.read()
    prs = Presentation(file.file)
    full_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"

    return {"conteudoExtraido": full_text.strip()}